import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def _blank_slide(prs: Presentation):
    # Blank layout
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_line_chart(prs: Presentation):
    slide = _blank_slide(prs)
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("S1", (1, 2, 3))
    x, y, cx, cy = 1_000_000, 1_000_000, 6_000_000, 4_000_000
    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    )
    return shape.chart


def _first_chart(slide):
    for shp in slide.shapes:
        if hasattr(shp, "chart"):
            return shp.chart
    raise AssertionError("No chart found on slide")


def test_datalabels_toggle_and_number_format_roundtrip():
    prs = Presentation()
    chart = _add_line_chart(prs)
    plot = chart.plots[0]

    plot.has_data_labels = True
    dlabels = plot.data_labels
    dlabels.show_value = True
    dlabels.show_category_name = False
    dlabels.show_series_name = False
    dlabels.number_format = "#,##0.00"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)

    chart2 = _first_chart(prs2.slides[0])
    d2 = chart2.plots[0].data_labels

    assert d2.show_value is True
    assert (d2.show_category_name or False) is False
    assert (d2.show_series_name or False) is False
    assert d2.number_format in ("#,##0.00",)


def test_line_chart_defaults_legend_and_vary_by_categories():
    prs = Presentation()
    slide = _blank_slide(prs)
    data = CategoryChartData()
    data.categories = ["Q1", "Q2"]
    data.add_series("S1", (1, 2))
    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, 0, 0, 6_000_000, 4_000_000, data
    )
    chart = shape.chart

    assert chart.has_legend is True
    assert chart.plots[0].vary_by_categories is False
