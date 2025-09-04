import io
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def _blank_slide(prs: Presentation):
    # Blank layout
    return prs.slides.add_slide(prs.slide_layouts[6])


def _roundtrip(prs: Presentation) -> Presentation:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _last_shape(slide):
    return slide.shapes[-1]


def test_hyperlink_http_roundtrip():
    prs = Presentation()
    slide = _blank_slide(prs)

    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1_000_000, 1_000_000, 2_000_000, 1_000_000
    )
    shp.click_action.hyperlink.address = "https://example.com/page#section"

    prs2 = _roundtrip(prs)
    shp2 = _last_shape(prs2.slides[0])
    addr = shp2.click_action.hyperlink.address or ""

    # Keep it simple and short to avoid E501
    assert "example.com/page#section" in addr


def test_hyperlink_file_uri_roundtrip(tmp_path: Path):
    prs = Presentation()
    slide = _blank_slide(prs)

    file_path = tmp_path / "doc.txt"
    file_path.write_text("hello")
    uri = file_path.as_uri()

    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0, 0, 2_000_000, 1_000_000
    )
    shp.click_action.hyperlink.address = uri

    prs2 = _roundtrip(prs)
    shp2 = _last_shape(prs2.slides[0])
    addr = shp2.click_action.hyperlink.address or ""

    assert addr.startswith("file:")
    assert file_path.name in addr


def test_hyperlink_mailto_roundtrip():
    prs = Presentation()
    slide = _blank_slide(prs)

    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0, 0, 2_000_000, 1_000_000
    )
    shp.click_action.hyperlink.address = "mailto:foo@example.com?subject=Hello"

    prs2 = _roundtrip(prs)
    shp2 = _last_shape(prs2.slides[0])
    addr = shp2.click_action.hyperlink.address or ""

    # PT018: break assertions
    assert addr.startswith("mailto:")
    assert "subject=Hello" in addr


def test_hyperlink_slide_anchor_via_target_slide_roundtrip():
    prs = Presentation()
    slide1 = _blank_slide(prs)
    slide2 = _blank_slide(prs)

    shp = slide1.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0, 0, 2_000_000, 1_000_000
    )
    # was: shp.click_action.hyperlink.target_slide = slide2
    shp.click_action.target_slide = slide2

    prs2 = _roundtrip(prs)
    shp2 = _last_shape(prs2.slides[0])
    # was: target = shp2.click_action.hyperlink.target_slide
    target = shp2.click_action.target_slide

    assert target is prs2.slides[1]

